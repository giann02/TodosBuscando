from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "Final Todos Buscando Presentacion - con flujo.pptx"

BG = RGBColor(15, 15, 30)
CARD = RGBColor(26, 26, 46)
TEAL = RGBColor(0, 212, 170)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(185, 190, 205)
BLUE = RGBColor(72, 149, 239)
YELLOW = RGBColor(255, 199, 95)
PURPLE = RGBColor(160, 120, 255)


def set_text(shape, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_textbox(slide, x, y, w, h, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_text(box, text, size=size, color=color, bold=bold, align=align)


def add_card(slide, x, y, w, h, num, title, body, accent=TEAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = accent
    card.line.width = Pt(1.25)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.16), Inches(y + 0.17), Inches(0.38), Inches(0.38))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    set_text(circle, str(num), size=11, color=BG, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide, x + 0.64, y + 0.17, w - 0.82, 0.28, title, size=11.5, color=WHITE, bold=True)
    add_textbox(slide, x + 0.18, y + 0.62, w - 0.36, h - 0.77, body, size=9.8, color=MUTED)


def add_arrow(slide, x1, y1, x2, y2, color=TEAL):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.8)
    line.line.end_arrowhead = True


def add_store(slide, x, y, w, label, title, body, color):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.76))
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD
    pill.line.color.rgb = color
    pill.line.width = Pt(1)
    add_textbox(slide, x + 0.18, y + 0.1, w - 0.36, 0.16, label, size=6.8, color=color, bold=True)
    add_textbox(slide, x + 0.18, y + 0.3, w - 0.36, 0.2, title, size=11, color=WHITE, bold=True)
    add_textbox(slide, x + 0.18, y + 0.54, w - 0.36, 0.14, body, size=6.8, color=MUTED)


def move_slide(prs, old_index: int, new_index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    sld_id_lst.remove(slides[old_index])
    sld_id_lst.insert(new_index, slides[old_index])


def main() -> None:
    prs = Presentation(PPTX_PATH)
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    # Motif used throughout the deck: small teal accent blocks.
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    add_textbox(slide, 0.7, 0.4, 10.0, 0.25, "FLUJO DE LA APLICACIÓN", size=10, color=TEAL, bold=True)
    add_textbox(slide, 0.7, 0.88, 11.2, 0.85, "De la alerta al rastro de avistamientos", size=34, color=WHITE, bold=True)
    add_textbox(
        slide,
        0.72,
        1.78,
        11.8,
        0.28,
        "El sistema combina entrada manual/API, consultas geoespaciales, notificación comunitaria y grafo de trayectoria.",
        size=12,
        color=MUTED,
    )

    cards = [
        ("Admin / API", "Se crea una alerta con datos, foto y última ubicación conocida.", TEAL),
        ("MongoDB", "Guarda la alerta y usa 2dsphere para ubicar vecinos cercanos.", BLUE),
        ("Email", "Vecinos dentro de 2 km reciben la notificación automáticamente.", YELLOW),
        ("Reporte", "Un vecino carga un avistamiento con descripción y ubicación opcional.", TEAL),
        ("Trayectoria", "MongoDB conserva el reporte y Neo4j conecta los puntos en orden.", PURPLE),
    ]

    x0, y0, w, h, gap = 0.55, 2.55, 2.28, 1.42, 0.27
    for index, (title, body, color) in enumerate(cards, start=1):
        x = x0 + (index - 1) * (w + gap)
        add_card(slide, x, y0, w, h, index, title, body, color)
        if index < len(cards):
            add_arrow(slide, x + w + 0.03, y0 + h / 2, x + w + gap - 0.05, y0 + h / 2, color=TEAL)

    add_textbox(slide, 0.72, 4.48, 3.2, 0.22, "CAPA DE DATOS INVOLUCRADA", size=8.5, color=TEAL, bold=True)
    add_store(slide, 0.72, 4.86, 3.55, "DOCUMENTAL", "MongoDB", "alertas, usuarios, reportes + geoespacial", BLUE)
    add_store(slide, 4.88, 4.86, 3.55, "CLAVE-VALOR", "Redis", "alertas activas con TTL e invalidación", YELLOW)
    add_store(slide, 9.04, 4.86, 3.55, "GRAFOS", "Neo4j", "cadena de avistamientos SIGUIENTE", PURPLE)

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(6.15), Inches(11.87), Inches(0.46))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(19, 32, 48)
    note.line.color.rgb = TEAL
    note.line.width = Pt(0.75)
    add_textbox(
        slide,
        0.95,
        6.28,
        11.4,
        0.18,
        "Resultado: el administrador ve alertas activas, reportes recibidos y el recorrido probable sobre el mapa.",
        size=10.5,
        color=WHITE,
        bold=False,
        align=PP_ALIGN.CENTER,
    )

    # Place this new flow slide before the final "CASOS DE USO" closing slide.
    move_slide(prs, len(prs.slides) - 1, len(prs.slides) - 2)
    prs.save(PPTX_PATH)


if __name__ == "__main__":
    main()
